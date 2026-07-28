# -*- coding: utf-8 -*-
"""집찾Go 부동산 웹 플랫폼 — 20분 발표용 10장 PPT 생성"""
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

REPO_ROOT = Path(__file__).resolve().parent.parent
OUT = REPO_ROOT / "집찾Go_부동산플랫폼_제작결과발표.pptx"

# Colors — real-estate / trust theme
NAVY = RGBColor(0x1A, 0x36, 0x5D)
TEAL = RGBColor(0x0D, 0x94, 0x88)
LIGHT = RGBColor(0xF4, 0xF7, 0xFA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x5A, 0x6A, 0x7A)
ACCENT = RGBColor(0x2E, 0x86, 0xAB)

FONT = "맑은 고딕"


def set_run(run, size=18, bold=False, color=None):
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color


def add_title_bar(slide, title_text, subtitle=None):
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.15)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()
    tf = bar.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = title_text
    set_run(r, 28, True, WHITE)
    tf.margin_left = Inches(0.5)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    if subtitle:
        box = slide.shapes.add_textbox(Inches(0.5), Inches(1.35), Inches(9), Inches(0.4))
        p2 = box.text_frame.paragraphs[0]
        r2 = p2.add_run()
        r2.text = subtitle
        set_run(r2, 14, False, GRAY)


def bullet_slide(prs, title, bullets, subtitle=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, title, subtitle)
    body = slide.shapes.add_textbox(Inches(0.55), Inches(1.85), Inches(8.9), Inches(5.2))
    tf = body.text_frame
    tf.word_wrap = True
    for i, item in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = 0
        p.space_after = Pt(10)
        r = p.add_run()
        r.text = item if item.startswith("•") else f"• {item}"
        set_run(r, 17, False, RGBColor(0x33, 0x33, 0x33))


def two_column_slide(prs, title, left_title, left_items, right_title, right_items):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, title)
    for col, (head, items, x) in enumerate(
        [(left_title, left_items, 0.5), (right_title, right_items, 5.2)]
    ):
        hbox = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.85), Inches(4.4), Inches(0.55)
        )
        hbox.fill.solid()
        hbox.fill.fore_color.rgb = TEAL if col == 0 else ACCENT
        hbox.line.fill.background()
        ht = hbox.text_frame
        ht.vertical_anchor = MSO_ANCHOR.MIDDLE
        hr = ht.paragraphs[0].add_run()
        hr.text = head
        set_run(hr, 16, True, WHITE)
        ht.paragraphs[0].alignment = PP_ALIGN.CENTER
        box = slide.shapes.add_textbox(Inches(x + 0.15), Inches(2.55), Inches(4.1), Inches(4.5))
        tf = box.text_frame
        for i, line in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            r = p.add_run()
            r.text = line
            set_run(r, 14, False, RGBColor(0x33, 0x33, 0x33))
            p.space_after = Pt(6)


def team_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "팀 구성원 소개", "4인 개발팀 · 집찾Go Real Estate Platform")
    members = [
        ("맹준형", "팀원", "기획 · 프론트/백엔드 개발"),
        ("이승우", "팀원", "UI/UX · 프론트엔드 · API 연동"),
        ("안종범", "팀원", "백엔드 · DB · 보안(Spring Security)"),
        ("황상옥", "팀원", "데이터 · RPA/AI · 시각화"),
    ]
    xs = [0.45, 2.55, 4.65, 6.75]
    for (name, role, desc), x in zip(members, xs):
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(2.0), Inches(2.0), Inches(3.8)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = LIGHT
        card.line.color.rgb = TEAL
        tf = card.text_frame
        tf.margin_top = Inches(0.25)
        tf.word_wrap = True
        p0 = tf.paragraphs[0]
        p0.alignment = PP_ALIGN.CENTER
        r0 = p0.add_run()
        r0.text = name
        set_run(r0, 20, True, NAVY)
        p1 = tf.add_paragraph()
        p1.alignment = PP_ALIGN.CENTER
        r1 = p1.add_run()
        r1.text = role
        set_run(r1, 12, False, TEAL)
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        p2.space_before = Pt(12)
        r2 = p2.add_run()
        r2.text = desc
        set_run(r2, 11, False, GRAY)
    note = slide.shapes.add_textbox(Inches(0.5), Inches(6.0), Inches(9), Inches(0.8))
    pn = note.text_frame.paragraphs[0]
    rn = pn.add_run()
    rn.text = "기획·디자인·개발·데이터 분석을 분담하여 사용자 중심 부동산 검색 AI 서비스를 구현했습니다."
    set_run(rn, 13, False, GRAY)


def title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(7.5)
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()
    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(4.2), Inches(10), Inches(0.08)
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = TEAL
    accent.line.fill.background()
    tbox = slide.shapes.add_textbox(Inches(0.6), Inches(1.8), Inches(8.8), Inches(2.2))
    tp = tbox.text_frame.paragraphs[0]
    tr = tp.add_run()
    tr.text = "집찾Go"
    set_run(tr, 54, True, WHITE)
    t2 = tbox.text_frame.add_paragraph()
    r2 = t2.add_run()
    r2.text = "부동산 웹 플랫폼 제작 결과 발표"
    set_run(r2, 32, True, RGBColor(0xB8, 0xE0, 0xDC))
    sub = slide.shapes.add_textbox(Inches(0.6), Inches(4.5), Inches(8.8), Inches(1.5))
    for i, line in enumerate(
        [
            "스마트 부동산 플랫폼 · 실거래·시장동향·AI 검색",
            "발표 시간: 약 20분 (10장)",
            "팀: 맹준형 · 이승우 · 안종범 · 황상옥",
        ]
    ):
        p = sub.text_frame.paragraphs[0] if i == 0 else sub.text_frame.add_paragraph()
        r = p.add_run()
        r.text = line
        set_run(r, 16, False, RGBColor(0xD0, 0xDC, 0xE8))


def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    title_slide(prs)

    bullet_slide(
        prs,
        "목차 (Agenda)",
        [
            "프로젝트 개요 및 목표",
            "팀 구성원 소개",
            "주요 기능 및 서비스",
            "시스템 아키텍처",
            "기술 스택 (Frontend / Backend)",
            "빌드·배포·Spring 의존성",
            "RPA + AI (Python 데이터 파이프라인)",
            "제작 결과 및 향후 계획",
        ],
        "슬라이드당 약 2분 · 총 10장",
    )

    bullet_slide(
        prs,
        "프로젝트 개요",
        [
            "서비스명: 집찾Go — 내 집 찾기의 시작",
            "목표: 사용자 중심 AI 부동산 검색·실거래 정보 제공",
            "핵심 가치: 쉬운 검색 UX, 실시간 매물·시장 정보, AI 맞춤 추천",
            "주요 메뉴: 실거래가, 시장동향, 부동산 뉴스, 매물 등록·중개 연계",
            "차별점: 웹(MVC) + 데이터 수집(RPA) + 분석·시각화(AI/ML) 통합",
        ],
        "ABOUT PROJECT",
    )

    team_slide(prs)

    bullet_slide(
        prs,
        "주요 기능 및 서비스",
        [
            "모바일 모델하우스 — 언제 어디서나 매물·단지 미리보기",
            "집 내놓기 — 간편 매물 등록 및 구매자 연결",
            "중개 라이브 — 실시간 중개·현장 확인 UX (확장 설계)",
            "회원·인증 — Spring Security + BCrypt 기반 로그인/권한",
            "대시보드 — 지역·가격대별 통계, Folium/Plotly 지도·차트 연동",
        ],
        "OUR SERVICES",
    )

    bullet_slide(
        prs,
        "시스템 아키텍처",
        [
            "Client: HTML5 · CSS3 · JavaScript · Bootstrap 5 · Flexbox",
            "Application: Spring Boot 3.5.x (JDK 21) · Thymeleaf SSR · Tomcat 10.1",
            "Data: MySQL 8 (로컬/개발) · TiDB (클라우드 확장·분산 DB)",
            "Build: Gradle — 의존성 관리 및 bootWar/JAR 패키징",
            "Deploy: Render 또는 AWS EC2 — CI 후 무중단 배포 목표",
            "Data Pipeline: Python RPA/크롤링 → 정제(Pandas) → DB/API 적재 → AI 분석",
        ],
    )

    two_column_slide(
        prs,
        "기술 스택 — Frontend & Backend",
        "Frontend",
        [
            "HTML / CSS / JavaScript",
            "Bootstrap 5",
            "Flexbox 레이아웃",
            "반응형 헤더·푸터 컴포넌트",
            "Thymeleaf 템플릿 연동",
        ],
        "Backend",
        [
            "JDK 21",
            "Spring Boot 3.5.x",
            "Thymeleaf (View)",
            "Embedded Tomcat 10.1",
            "MySQL 8 + JDBC Template",
        ],
    )

    two_column_slide(
        prs,
        "빌드 · 배포 · Spring 의존성",
        "빌드 & 배포",
        [
            "Gradle 빌드",
            "TiDB (클라우드 DB)",
            "Render 또는 EC2 배포",
            "환경별 application.yml",
            "HTTPS·도메인 연동",
        ],
        "주요 의존성 (DI)",
        [
            "spring-boot-starter-web",
            "spring-boot-starter-security (BCrypt)",
            "spring-boot-starter-thymeleaf",
            "mysql-connector-j",
            "Lombok",
            "spring-boot-starter-jdbc (JdbcTemplate)",
        ],
    )

    two_column_slide(
        prs,
        "RPA + AI — Python 패키지",
        "데이터 수집·전처리",
        [
            "Requests — HTTP API·페이지 요청",
            "BeautifulSoup — HTML 파싱(RPA/크롤링)",
            "NumPy · Pandas — 배열·표 데이터 처리",
            "Scikit-learn — 추천·분류·회귀 모델",
        ],
        "시각화·지도",
        [
            "Matplotlib · Seaborn — 통계 차트",
            "Plotly — 인터랙티브 그래프",
            "Folium — 지역 매물·실거래 지도",
            "→ Spring API/DB와 연동해 대시보드 제공",
        ],
    )

    bullet_slide(
        prs,
        "제작 결과 및 마무리",
        [
            "완료: 부동산 플랫폼 UI/UX, Spring Boot 기반 서버·보안·DB 연동 골격",
            "완료: Gradle 빌드 파이프라인, 팀 소개·랜딩 페이지, 데이터 분석 PoC",
            "성과: Full-stack + 데이터/AI 역량을 하나의 서비스로 통합",
            "향후: TiDB·Render/EC2 운영 배포, AI 추천 고도화, 실거래 API 정식 연동",
            "감사합니다 · Q & A",
        ],
        "Thank you",
    )

    prs.save(str(OUT))
    print(f"Saved: {OUT}")


if __name__ == "__main__":
    main()
